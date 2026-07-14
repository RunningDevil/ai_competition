from __future__ import annotations

import posixpath
import re
import zipfile
from pathlib import Path
from typing import Any, Dict, List
from xml.etree import ElementTree as ET

from comment_parser import looks_like_comment, parse_comment_text
from office_common import (
    convert_legacy_office,
    extract_printable_text_from_binary,
    filter_comments,
    fixed_relative_path,
    get_extension,
    make_error,
    normalize_path_text,
    replacement_from_comment,
    source_relative_to_docs,
)


A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
RELS_NS = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
A_P = f"{{{A_NS['a']}}}p"
A_T = f"{{{A_NS['a']}}}t"


def _extract_text_from_shape(shape: Any) -> str:
    if not hasattr(shape, "text"):
        return ""
    return str(shape.text or "").strip()


def _extract_pptx_text(path: Path, source_rel: str) -> List[Dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    texts: List[Dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = _extract_text_from_shape(shape)
            if text:
                texts.append({
                    "source": source_rel,
                    "file_type": "pptx",
                    "location": f"slide:{slide_index}:shape:{shape_index}",
                    "text": text,
                })
    return texts


def _xml_text(package: zipfile.ZipFile, name: str) -> str:
    try:
        root = ET.fromstring(package.read(name))
    except Exception:
        return ""
    parts = [node.text or "" for node in root.findall(".//a:t", A_NS)]
    return "\n".join(part.strip() for part in parts if part and part.strip())


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def _normalize_part_path(path: str) -> str:
    return posixpath.normpath(path).lstrip("/")


def _target_part_from_rels(rels_name: str, target: str) -> str:
    # ppt/slides/_rels/slide1.xml.rels targets are relative to ppt/slides/.
    source_part = rels_name.replace("/_rels/", "/")
    if source_part.endswith(".rels"):
        source_part = source_part[: -len(".rels")]
    base_dir = posixpath.dirname(source_part)
    return _normalize_part_path(posixpath.join(base_dir, target))


def _comment_authors(package: zipfile.ZipFile) -> Dict[str, str]:
    if "ppt/commentAuthors.xml" not in package.namelist():
        return {}
    try:
        root = ET.fromstring(package.read("ppt/commentAuthors.xml"))
    except Exception:
        return {}
    authors: Dict[str, str] = {}
    for node in root.iter():
        if _local_name(node.tag) != "cmAuthor":
            continue
        author_id = node.attrib.get("id")
        name = node.attrib.get("name") or node.attrib.get("initials") or ""
        if author_id is not None and name:
            authors[str(author_id)] = name
    return authors


def _comment_slide_map(package: zipfile.ZipFile) -> Dict[str, int]:
    mapping: Dict[str, int] = {}
    for rels_name in package.namelist():
        match = re.fullmatch(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", rels_name)
        if not match:
            continue
        slide_index = int(match.group(1))
        try:
            root = ET.fromstring(package.read(rels_name))
        except Exception:
            continue
        for rel in root.findall("rel:Relationship", RELS_NS):
            rel_type = rel.attrib.get("Type", "").lower()
            target = rel.attrib.get("Target", "")
            if not target or "comment" not in rel_type:
                continue
            mapping[_target_part_from_rels(rels_name, target)] = slide_index
    return mapping


def _comment_part_names(package: zipfile.ZipFile) -> List[str]:
    names = []
    for name in package.namelist():
        lowered = name.lower()
        if not lowered.endswith(".xml"):
            continue
        if lowered.startswith("ppt/comments/comment") or lowered.startswith("ppt/threadedcomments/threadedcomment"):
            names.append(name)
    return sorted(names)


def _text_from_comment_element(element: ET.Element) -> str:
    parts: List[str] = []
    for node in element.iter():
        local = _local_name(node.tag)
        if local in {"text", "t"} and node.text:
            text = node.text.strip()
            if text:
                parts.append(text)
    if not parts:
        parts = [text.strip() for text in element.itertext() if text and text.strip()]
    return "\n".join(parts).strip()


def _comment_elements(root: ET.Element) -> List[ET.Element]:
    candidates = [node for node in root.iter() if _local_name(node.tag) in {"cm", "comment", "threadedComment"}]
    if candidates:
        return candidates
    text = _text_from_comment_element(root)
    return [root] if text else []


def _extract_pptx_comments(path: Path, source_rel: str) -> List[Dict[str, Any]]:
    comments: List[Dict[str, Any]] = []
    with zipfile.ZipFile(path) as package:
        authors = _comment_authors(package)
        slide_map = _comment_slide_map(package)
        for name in _comment_part_names(package):
            try:
                root = ET.fromstring(package.read(name))
            except Exception:
                continue
            slide_index = slide_map.get(name)
            for index, node in enumerate(_comment_elements(root), start=1):
                text = _text_from_comment_element(node)
                if not text:
                    continue
                location_prefix = f"slide:{slide_index}:" if slide_index else ""
                comment = parse_comment_text(text, source_rel, "pptx", f"{location_prefix}{name}:comment:{index}")
                author_id = node.attrib.get("authorId") or node.attrib.get("author")
                if author_id is not None:
                    comment["author_id"] = str(author_id)
                    if str(author_id) in authors:
                        comment["author"] = authors[str(author_id)]
                if node.attrib.get("dt"):
                    comment["created_at"] = node.attrib.get("dt")
                comments.append(comment)
    return comments


def _replace_in_slide_xml(xml_bytes: bytes, old: str, new: str) -> tuple[bytes, int]:
    root = ET.fromstring(xml_bytes)
    replaced_count = 0
    for paragraph in root.iter(A_P):
        text_nodes = [node for node in paragraph.iter(A_T)]
        if not text_nodes:
            continue
        combined = "".join(node.text or "" for node in text_nodes)
        if old not in combined:
            continue
        updated = combined.replace(old, new, 1)
        text_nodes[0].text = updated
        for node in text_nodes[1:]:
            node.text = ""
        replaced_count += 1
        break
    return ET.tostring(root, encoding="utf-8", xml_declaration=True), replaced_count


def _write_replaced_pptx(source: Path, target: Path, old: str, new: str, logs: List[str]) -> int:
    target.parent.mkdir(parents=True, exist_ok=True)
    replaced_count = 0
    slide_name_re = re.compile(r"ppt/slides/slide\d+\.xml")
    with zipfile.ZipFile(source, "r") as zin, zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if replaced_count == 0 and slide_name_re.fullmatch(item.filename):
                data, replaced_count = _replace_in_slide_xml(data, old, new)
            zout.writestr(item, data)
    if replaced_count:
        logs.append(f"Replaced PowerPoint text: {old} -> {new}")
    else:
        logs.append(f"Replacement text not found in PowerPoint slides: {old}")
    return replaced_count


def extract_text(path: Path, wiki_root: str | Path, logs: List[str]) -> List[Dict[str, Any]]:
    ext = get_extension(path)
    source_rel = source_relative_to_docs(path, wiki_root)
    if ext == ".pptx":
        return _extract_pptx_text(path, source_rel)
    if ext == ".ppt":
        converted = convert_legacy_office(path, logs)
        if converted and converted.suffix.lower() == ".pptx":
            return _extract_pptx_text(converted, source_rel)
        fallback = extract_printable_text_from_binary(path)
        return [{"source": source_rel, "file_type": "ppt", "location": "legacy:fallback_text", "text": fallback}] if fallback else []
    return []


def extract_comments(path: Path, wiki_root: str | Path, logs: List[str]) -> List[Dict[str, Any]]:
    ext = get_extension(path)
    source_rel = source_relative_to_docs(path, wiki_root)
    if ext == ".pptx":
        return _extract_pptx_comments(path, source_rel)
    if ext == ".ppt":
        converted = convert_legacy_office(path, logs)
        if converted and converted.suffix.lower() == ".pptx":
            return _extract_pptx_comments(converted, source_rel)
        fallback = extract_printable_text_from_binary(path)
        if looks_like_comment(fallback):
            return [parse_comment_text(fallback, source_rel, "ppt", "legacy:fallback_text")]
        logs.append("Legacy .ppt comments could not be reliably extracted")
    return []


def conservative_fix(path: Path, wiki_root: str | Path, logs: List[str], filters: Dict[str, Any] | None = None) -> Dict[str, Any]:
    if get_extension(path) != ".pptx":
        return make_error("fix_comments", f"PowerPoint repair only supports .pptx for now: {normalize_path_text(str(path))}", logs)
    comments = extract_comments(path, wiki_root, logs)
    if filters:
        comments = filter_comments(comments, filters)
    if not comments:
        return make_error("fix_comments", f"No reliable PowerPoint comments found for {normalize_path_text(str(path))}", logs)
    source_rel = source_relative_to_docs(path, wiki_root)
    target_rel = fixed_relative_path(source_rel)
    target_abs = Path(wiki_root) / target_rel
    for comment in comments:
        replacement = replacement_from_comment(comment)
        if not replacement:
            continue
        old, new = replacement
        replaced_count = _write_replaced_pptx(path, target_abs, old, new, logs)
        if replaced_count:
            return {
                "status": "ok",
                "task_type": "fix_comments",
                "texts": [],
                "comments": comments,
                "answer": {"source": source_rel, "target": target_rel},
                "fixed_files": [target_rel],
                "logs": logs,
            }
    logs.append("PowerPoint comments found, but no deterministic replacement instruction could be applied")
    return make_error("fix_comments", f"No reliable PowerPoint repair rule implemented for {normalize_path_text(str(path))}", logs)
